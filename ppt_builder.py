from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

THEMES = {
    "elegant": {
        "bg1":           RGBColor(0x1A, 0x1A, 0x2E),
        "bg2":           RGBColor(0x16, 0x21, 0x3E),
        "accent":        RGBColor(0xE2, 0xC2, 0x7D),
        "text":          RGBColor(0xFF, 0xFF, 0xFF),
        "sub":           RGBColor(0xCC, 0xAA, 0x66),
        "bullet_bg":     RGBColor(0x0D, 0x11, 0x22),
        "bullet_border": RGBColor(0x2A, 0x35, 0x5A),
        "arrow":         RGBColor(0xE2, 0xC2, 0x7D),
    },
    "ocean": {
        "bg1":           RGBColor(0x0F, 0x20, 0x44),
        "bg2":           RGBColor(0x0A, 0x18, 0x36),
        "accent":        RGBColor(0x7D, 0xD3, 0xFC),
        "text":          RGBColor(0xFF, 0xFF, 0xFF),
        "sub":           RGBColor(0x93, 0xC5, 0xFD),
        "bullet_bg":     RGBColor(0x07, 0x12, 0x28),
        "bullet_border": RGBColor(0x1E, 0x40, 0x7A),
        "arrow":         RGBColor(0x38, 0xBD, 0xF8),
    },
    "rose": {
        "bg1":           RGBColor(0x1C, 0x0A, 0x14),
        "bg2":           RGBColor(0x2D, 0x0A, 0x1E),
        "accent":        RGBColor(0xFD, 0xA4, 0xAF),
        "text":          RGBColor(0xFF, 0xFF, 0xFF),
        "sub":           RGBColor(0xFB, 0xBC, 0xC5),
        "bullet_bg":     RGBColor(0x10, 0x05, 0x0C),
        "bullet_border": RGBColor(0x5E, 0x10, 0x30),
        "arrow":         RGBColor(0xE1, 0x1D, 0x6A),
    },
    "forest": {
        "bg1":           RGBColor(0x05, 0x2E, 0x16),
        "bg2":           RGBColor(0x03, 0x1E, 0x0E),
        "accent":        RGBColor(0xBB, 0xF7, 0xD0),
        "text":          RGBColor(0xFF, 0xFF, 0xFF),
        "sub":           RGBColor(0x86, 0xEF, 0xAC),
        "bullet_bg":     RGBColor(0x02, 0x10, 0x08),
        "bullet_border": RGBColor(0x14, 0x53, 0x2D),
        "arrow":         RGBColor(0x16, 0xA3, 0x4A),
    },
    "minimal": {
        "bg1":           RGBColor(0xFA, 0xFA, 0xFA),
        "bg2":           RGBColor(0xF3, 0xF4, 0xF6),
        "accent":        RGBColor(0x11, 0x11, 0x11),
        "text":          RGBColor(0x11, 0x11, 0x11),
        "sub":           RGBColor(0x6B, 0x72, 0x80),
        "bullet_bg":     RGBColor(0xFF, 0xFF, 0xFF),
        "bullet_border": RGBColor(0xD1, 0xD5, 0xDB),
        "arrow":         RGBColor(0x9C, 0xA3, 0xAF),
    },
    "sunset": {
        "bg1":           RGBColor(0x1C, 0x0A, 0x00),
        "bg2":           RGBColor(0x2D, 0x10, 0x00),
        "accent":        RGBColor(0xFC, 0xD3, 0x4D),
        "text":          RGBColor(0xFF, 0xFF, 0xFF),
        "sub":           RGBColor(0xFD, 0xBA, 0x74),
        "bullet_bg":     RGBColor(0x10, 0x05, 0x00),
        "bullet_border": RGBColor(0x7C, 0x2D, 0x12),
        "arrow":         RGBColor(0xEA, 0x58, 0x0C),
    },
    "violet": {
        "bg1":           RGBColor(0x0D, 0x07, 0x20),
        "bg2":           RGBColor(0x13, 0x0A, 0x30),
        "accent":        RGBColor(0xC4, 0xB5, 0xFD),
        "text":          RGBColor(0xFF, 0xFF, 0xFF),
        "sub":           RGBColor(0xA7, 0x8B, 0xFA),
        "bullet_bg":     RGBColor(0x08, 0x04, 0x14),
        "bullet_border": RGBColor(0x3B, 0x1A, 0x7A),
        "arrow":         RGBColor(0x7C, 0x3A, 0xED),
    },
    "slate": {
        "bg1":           RGBColor(0x0F, 0x17, 0x2A),
        "bg2":           RGBColor(0x1E, 0x29, 0x3B),
        "accent":        RGBColor(0xE2, 0xE8, 0xF0),
        "text":          RGBColor(0xFF, 0xFF, 0xFF),
        "sub":           RGBColor(0x94, 0xA3, 0xB8),
        "bullet_bg":     RGBColor(0x0A, 0x10, 0x1E),
        "bullet_border": RGBColor(0x33, 0x44, 0x55),
        "arrow":         RGBColor(0x64, 0x74, 0x8B),
    },
}


def add_textbox(slide, text, left, top, width, height,
                size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return box


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def make_title_slide(prs, title, subtitle, t):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, t["bg1"])

    # Left accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = t["accent"]
    bar.line.fill.background()

    # Bottom strip
    strip = slide.shapes.add_shape(1, Inches(0), Inches(7.05), Inches(13.33), Inches(0.45))
    strip.fill.solid()
    strip.fill.fore_color.rgb = t["bg2"]
    strip.line.fill.background()

    # Decorative top-right corner box
    corner = slide.shapes.add_shape(1, Inches(11.5), Inches(0), Inches(1.83), Inches(1.5))
    corner.fill.solid()
    corner.fill.fore_color.rgb = t["bullet_bg"]
    corner.line.fill.background()

    add_textbox(slide, title,    0.65, 2.1, 11.5, 2.2, size=44, bold=True,  color=t["text"])
    add_textbox(slide, subtitle, 0.65, 4.5, 10.5, 1.0, size=20, bold=False, color=t["sub"])
    add_textbox(slide, "Generated by SlideAI  ·  Powered by Groq",
                0.65, 6.8, 8.0, 0.4, size=11, bold=False, color=t["sub"])


def make_content_slide(prs, heading, bullets, speaker_note, t):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, t["bg2"])

    # Top accent bar
    top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.12))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = t["accent"]
    top_bar.line.fill.background()

    # Heading background
    head_bg = slide.shapes.add_shape(1, Inches(0.3), Inches(0.18), Inches(12.6), Inches(1.0))
    head_bg.fill.solid()
    head_bg.fill.fore_color.rgb = t["bg1"]
    head_bg.line.fill.background()

    add_textbox(slide, heading, 0.5, 0.2, 12.0, 0.95, size=28, bold=True, color=t["accent"])

    # Bullets
    for i, bullet in enumerate(bullets):
        top = 1.42 + i * 0.92

        bg = slide.shapes.add_shape(1, Inches(0.3), Inches(top), Inches(12.6), Inches(0.78))
        bg.fill.solid()
        bg.fill.fore_color.rgb = t["bullet_bg"]
        bg.line.color.rgb = t["bullet_border"]
        bg.line.width = Pt(0.75)

        add_textbox(slide, "▸", 0.42, top + 0.12, 0.38, 0.56, size=13, bold=True, color=t["arrow"])
        add_textbox(slide, bullet, 0.84, top + 0.1, 12.0, 0.62, size=16, bold=False, color=t["text"])

    # Speaker note
    if speaker_note:
        slide.notes_slide.notes_text_frame.text = speaker_note


def build_ppt(data: dict, filename: str = "output.pptx", theme: str = "elegant"):
    t = THEMES.get(theme, THEMES["elegant"])

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs, data["title"], data["subtitle"], t)
    print(f"[builder] Title slide: {data['title']}  |  Theme: {theme}")

    for i, s in enumerate(data["slides"]):
        make_content_slide(
            prs,
            heading=s["heading"],
            bullets=s["bullets"],
            speaker_note=s.get("speaker_note", ""),
            t=t
        )
        print(f"[builder] Slide {i+1}: {s['heading']}")

    prs.save(filename)
    print(f"[builder] Saved → {filename}")


# Quick test
if __name__ == "__main__":
    sample = {
        "title": "The Future of AI",
        "subtitle": "How artificial intelligence is reshaping our world",
        "slides": [
            {"heading": "Introduction", "bullets": ["AI is transforming industries", "Rapid growth since 2020", "Billions in investment globally"], "speaker_note": "Welcome everyone."},
            {"heading": "Key Technologies", "bullets": ["Large language models", "Computer vision advances", "Reinforcement learning"], "speaker_note": "Let's dive into tech."},
            {"heading": "Conclusion", "bullets": ["AI is here to stay", "Adapt and embrace change", "Exciting future ahead"], "speaker_note": "Thank you."},
        ]
    }
    for theme in ["elegant", "ocean", "rose", "violet"]:
        build_ppt(sample, f"test_{theme}.pptx", theme)
        print(f"Created test_{theme}.pptx")
